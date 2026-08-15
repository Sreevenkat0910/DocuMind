"""DocuMind Embedding & LLM Gateway (Plan.md Phase 3).

Smallest possible first cut of the Python service: text in -> embedding vector out,
and prompt in -> LLM completion out. No vector store, no retrieval, no permissions —
Java still owns all of that. Every route requires the shared internal API key.
"""
import json
import logging
import io
import os
import re
import time
import uuid
from contextlib import asynccontextmanager

from fastapi import Depends, FastAPI, Header, HTTPException, Response, UploadFile
from fastapi.responses import StreamingResponse
from groq import Groq
from langchain_text_splitters import RecursiveCharacterTextSplitter
from psycopg.types.json import Json
from psycopg_pool import ConnectionPool
from pydantic import BaseModel
from sentence_transformers import CrossEncoder, SentenceTransformer

API_KEY = os.environ["INTERNAL_API_KEY"]  # required: fail fast if unset
EMBED_MODEL = os.getenv("EMBED_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
RERANK_MODEL = os.getenv("RERANK_MODEL", "cross-encoder/ms-marco-MiniLM-L-6-v2")
GROQ_MODEL = os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile")
DATABASE_URL = os.getenv("DATABASE_URL", "postgresql://postgres:postgres@localhost:5432/documind")

# Structure-aware chunking (Plan.md Phase 6). Character-based sizing (~5 chars/word), roughly
# equivalent to the old 200-word/40-word Java sliding window, but the splitter respects
# paragraph/sentence boundaries instead of blindly cutting mid-word every N tokens.
CHUNK_SIZE_CHARS = int(os.getenv("CHUNK_SIZE_CHARS", "1000"))
CHUNK_OVERLAP_CHARS = int(os.getenv("CHUNK_OVERLAP_CHARS", "200"))

# Hybrid search (Plan.md Phase 6): how many candidates each of vector search and full-text
# search contributes before RRF fusion + reranking narrow down to the caller's top_k.
CANDIDATE_POOL = int(os.getenv("HYBRID_CANDIDATE_POOL", "20"))

# Grounding rules, ported verbatim from Java's ChatService (Plan.md Phase 5). Python now
# owns prompt construction, so this is the single source of truth for how the LLM answers.
SYSTEM_INSTRUCTION = (
    "You are Docent, an enterprise knowledge assistant. Answer the user's question using "
    "ONLY the information in the provided context. If the answer is not contained in the "
    "context, say you don't have enough information to answer — do not use outside knowledge "
    "and do not guess. Be concise, and refer to the bracketed source numbers (e.g. [1]) you used."
)

_embedder = SentenceTransformer(EMBED_MODEL)  # loaded once at startup (384-dim)
_reranker = CrossEncoder(RERANK_MODEL)  # loaded once at startup, same pattern as the embedder
_groq = Groq(api_key=os.getenv("GROQ_API_KEY") or "gsk-not-set")
# Opened lazily on the first DB use so importing this module (tests, /embed-only use)
# needs no database.
_pool = ConnectionPool(DATABASE_URL, min_size=1, max_size=4, open=False)

# --- Phase 7: structured, per-stage tracing. One JSON line per stage, always carrying the
# correlation id Java propagated (CorrelationIdFilter -> GatewayClient -> X-Correlation-Id
# header) — `grep <id>` across both services' logs shows one request's full timing story. No
# Langfuse/external tracing service: no account exists for this project, and stdlib logging
# already satisfies the actual exit criterion (see one trace, see where the latency went).
logging.basicConfig(level=logging.INFO, format="%(message)s")
_log = logging.getLogger("documind.gateway")


def log_stage(correlation_id: str, stage: str, duration_ms: float, **extra):
    """Emit one structured JSON log line for a named pipeline stage (retrieve / rerank / llm)."""
    payload = {"correlation_id": correlation_id, "stage": stage, "duration_ms": round(duration_ms, 1)}
    payload.update(extra)
    _log.info(json.dumps(payload))


def correlation_id_dep(x_correlation_id: str | None = Header(default=None)) -> str:
    """Reuse the caller's correlation id if it sent one (the normal case — Java always does),
    otherwise mint one so direct/manual calls still get a traceable id."""
    return x_correlation_id or str(uuid.uuid4())


def _ensure_schema():
    """Create the vector_store table + HNSW index Python now owns (Plan.md Phase 5.5 —
    this replaces Spring AI's initialize-schema). Idempotent: IF NOT EXISTS no-ops when the
    table already exists (e.g. the dev DB where Spring AI created it earlier)."""
    _pool.open()  # idempotent
    with _pool.connection() as conn:
        conn.execute("CREATE EXTENSION IF NOT EXISTS vector")
        conn.execute(
            "CREATE TABLE IF NOT EXISTS vector_store ("
            "  id uuid PRIMARY KEY,"
            "  content text,"
            "  metadata jsonb,"
            "  embedding vector(384))"
        )
        conn.execute(
            "CREATE INDEX IF NOT EXISTS vector_store_embedding_idx "
            "ON vector_store USING hnsw (embedding vector_cosine_ops)"
        )
        # Phase 6: full-text search side of hybrid retrieval. A generated column (auto-kept in
        # sync with `content`) + GIN index is the standard, boring way to do this in Postgres —
        # no separate BM25 index/library to maintain in Python.
        conn.execute(
            "ALTER TABLE vector_store ADD COLUMN IF NOT EXISTS content_tsv tsvector "
            "GENERATED ALWAYS AS (to_tsvector('english', coalesce(content, ''))) STORED"
        )
        conn.execute(
            "CREATE INDEX IF NOT EXISTS vector_store_content_tsv_idx "
            "ON vector_store USING gin (content_tsv)"
        )


@asynccontextmanager
async def _lifespan(_app: FastAPI):
    _ensure_schema()  # runs when the server boots; module-level TestClient (no-DB tests) skips it
    yield


app = FastAPI(title="DocuMind Embedding & LLM Gateway", lifespan=_lifespan)


def require_key(x_internal_api_key: str = Header(default="")):
    if x_internal_api_key != API_KEY:
        raise HTTPException(status_code=401, detail="bad internal api key")


class EmbedRequest(BaseModel):
    text: str


class EmbedResponse(BaseModel):
    embedding: list[float]
    model: str
    dim: int


class EmbedBatchRequest(BaseModel):
    texts: list[str]


class EmbedBatchResponse(BaseModel):
    embeddings: list[list[float]]
    model: str
    dim: int


class Message(BaseModel):
    role: str
    content: str


class GenerateRequest(BaseModel):
    system: str | None = None
    messages: list[Message]


class GenerateResponse(BaseModel):
    content: str
    model: str
    usage: dict


class RetrieveRequest(BaseModel):
    query: str
    allowed_document_ids: list[str]
    top_k: int = 5


class RetrieveResult(BaseModel):
    chunk_id: str | None
    document_id: str | None
    score: float


class RetrieveResponse(BaseModel):
    results: list[RetrieveResult]


@app.get("/health")
def health():
    return {"status": "ok", "embed_model": EMBED_MODEL,
            "dim": _embedder.get_sentence_embedding_dimension()}


@app.post("/embed", response_model=EmbedResponse, dependencies=[Depends(require_key)])
def embed(req: EmbedRequest):
    vec = _embedder.encode(req.text).tolist()
    return EmbedResponse(embedding=vec, model=EMBED_MODEL, dim=len(vec))


@app.post("/embed/batch", response_model=EmbedBatchResponse, dependencies=[Depends(require_key)])
def embed_batch(req: EmbedBatchRequest):
    vecs = [v.tolist() for v in _embedder.encode(req.texts)]
    dim = len(vecs[0]) if vecs else _embedder.get_sentence_embedding_dimension()
    return EmbedBatchResponse(embeddings=vecs, model=EMBED_MODEL, dim=dim)


# --- Phase 6: structure-aware chunking. Java still owns DocumentChunk persistence (it needs
# real JPA-generated chunk ids), so this only returns the split text + page number for each
# chunk in order; Java assigns chunkIndex and ids when it persists them.

class ChunkPage(BaseModel):
    page_number: int
    text: str


class ChunkRequest(BaseModel):
    pages: list[ChunkPage]


class ChunkResult(BaseModel):
    content: str
    page_number: int


class ChunkResponse(BaseModel):
    chunks: list[ChunkResult]


_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE_CHARS,
    chunk_overlap=CHUNK_OVERLAP_CHARS,
    separators=["\n\n", "\n", ". ", " ", ""],  # paragraph -> line -> sentence -> word -> char
)


def split_page(text: str, page_number: int) -> list[dict]:
    """Pure function: split one page's text into structure-aware chunks. Chunks never cross a
    page boundary (keeps `page_number` exact for citations, same invariant as the old Java
    sliding window). A blank page contributes no chunks."""
    stripped = text.strip() if text else ""
    if not stripped:
        return []
    return [{"content": c, "page_number": page_number} for c in _splitter.split_text(stripped)]


@app.post("/chunk", response_model=ChunkResponse, dependencies=[Depends(require_key)])
def chunk_document(req: ChunkRequest):
    chunks = []
    for page in req.pages:
        chunks.extend(split_page(page.text, page.page_number))
    return ChunkResponse(chunks=[ChunkResult(**c) for c in chunks])


@app.post("/generate", response_model=GenerateResponse, dependencies=[Depends(require_key)])
def generate(req: GenerateRequest):
    msgs = ([{"role": "system", "content": req.system}] if req.system else [])
    msgs += [{"role": m.role, "content": m.content} for m in req.messages]
    resp = _groq.chat.completions.create(model=GROQ_MODEL, messages=msgs)
    usage = resp.usage.model_dump() if resp.usage else {}
    return GenerateResponse(content=resp.choices[0].message.content, model=resp.model, usage=usage)


# --- Phase 6: hybrid search (vector + full-text) + CrossEncoder reranking. Shared by /retrieve
# and /rag/query so there's one retrieval implementation, not two. Java passes the already-
# computed allowed_document_ids and RE-VERIFIES every returned chunk_id — never an authorization
# check, only a filter + relevance ranking.

_VECTOR_CANDIDATES_SQL = (
    "SELECT metadata->>'chunkId' AS chunk_id, metadata->>'documentId' AS document_id, content, "
    "       1 - (embedding <=> %(vec)s::vector) AS score "
    "FROM vector_store "
    "WHERE metadata->>'documentId' = ANY(%(docs)s) "
    "ORDER BY embedding <=> %(vec)s::vector "
    "LIMIT %(n)s"
)

_LEXICAL_CANDIDATES_SQL = (
    "SELECT metadata->>'chunkId' AS chunk_id, metadata->>'documentId' AS document_id, content, "
    "       ts_rank_cd(content_tsv, plainto_tsquery('english', %(q)s)) AS score "
    "FROM vector_store "
    "WHERE metadata->>'documentId' = ANY(%(docs)s) "
    "  AND content_tsv @@ plainto_tsquery('english', %(q)s) "
    "ORDER BY score DESC "
    "LIMIT %(n)s"
)


def _row_to_hit(row) -> dict:
    return {"chunk_id": row[0], "document_id": row[1], "content": row[2], "score": float(row[3])}


def rrf_fuse(vector_hits: list[dict], lexical_hits: list[dict], k: int = 60) -> list[dict]:
    """Reciprocal Rank Fusion: combine two ranked candidate lists into one, ranked by
    sum(1 / (k + rank)) across whichever list(s) a chunk appears in. Pure function — no DB,
    no model — so it's directly unit-testable. Hits with no chunk_id (orphan metadata) are
    dropped rather than trusted."""
    scores: dict[str, float] = {}
    items: dict[str, dict] = {}
    for ranked_list in (vector_hits, lexical_hits):
        for rank, hit in enumerate(ranked_list):
            cid = hit.get("chunk_id")
            if not cid:
                continue
            scores[cid] = scores.get(cid, 0.0) + 1.0 / (k + rank + 1)
            items.setdefault(cid, hit)
    return sorted(items.values(), key=lambda h: scores[h["chunk_id"]], reverse=True)


def rerank(query: str, candidates: list[dict], top_k: int) -> list[dict]:
    """CrossEncoder reranking: score each (query, chunk content) pair, keep the top_k, and
    replace each hit's `score` with the CrossEncoder's relevance score (not the fused
    vector/lexical score it arrived with)."""
    if not candidates:
        return []
    pairs = [(query, c["content"]) for c in candidates]
    ce_scores = _reranker.predict(pairs)
    ranked = sorted(zip(candidates, ce_scores), key=lambda cs: cs[1], reverse=True)
    return [{**c, "score": float(s)} for c, s in ranked[:top_k]]


def hybrid_retrieve(query: str, allowed_document_ids: list[str], top_k: int,
                     correlation_id: str = "-") -> list[dict]:
    """Vector ANN + Postgres full-text search (top ~20 each) -> RRF fusion -> CrossEncoder
    rerank -> top_k. Returns hits as {chunk_id, document_id, content, score} dicts, `score`
    being the CrossEncoder's relevance score (not a similarity/distance metric any more).
    Logs the "retrieve" (search + fusion) and "rerank" stages (Plan.md Phase 7)."""
    if not allowed_document_ids:
        return []

    t0 = time.monotonic()
    vec = "[" + ",".join(map(str, _embedder.encode(query).tolist())) + "]"
    _pool.open()  # idempotent; establishes the pool on first use
    with _pool.connection() as conn:
        vector_rows = conn.execute(
            _VECTOR_CANDIDATES_SQL,
            {"vec": vec, "docs": allowed_document_ids, "n": CANDIDATE_POOL},
        ).fetchall()
        lexical_rows = conn.execute(
            _LEXICAL_CANDIDATES_SQL,
            {"q": query, "docs": allowed_document_ids, "n": CANDIDATE_POOL},
        ).fetchall()
    fused = rrf_fuse([_row_to_hit(r) for r in vector_rows], [_row_to_hit(r) for r in lexical_rows])
    log_stage(correlation_id, "retrieve", (time.monotonic() - t0) * 1000,
              vector_candidates=len(vector_rows), lexical_candidates=len(lexical_rows))

    t1 = time.monotonic()
    reranked = rerank(query, fused, top_k)
    log_stage(correlation_id, "rerank", (time.monotonic() - t1) * 1000,
              candidates_in=len(fused), results_out=len(reranked))
    return reranked


@app.post("/retrieve", response_model=RetrieveResponse, dependencies=[Depends(require_key)])
def retrieve(req: RetrieveRequest, response: Response, correlation_id: str = Depends(correlation_id_dep)):
    response.headers["X-Correlation-Id"] = correlation_id
    hits = hybrid_retrieve(req.query, req.allowed_document_ids, req.top_k, correlation_id)
    return RetrieveResponse(results=[
        RetrieveResult(chunk_id=h["chunk_id"], document_id=h["document_id"], score=h["score"])
        for h in hits
    ])


# --- Phase 5: full RAG assembly (retrieve -> prompt -> LLM -> cited sources) lives here now.
# Java passes the pre-computed allowed_document_ids and RE-VERIFIES every returned chunk_id
# against a fresh permission query — this endpoint is trusted to filter, never to authorize.

class RagQueryRequest(BaseModel):
    question: str
    allowed_document_ids: list[str]
    top_k: int = 5


class RagQueryResponse(BaseModel):
    answer: str
    used_chunk_ids: list[str]
    no_context_found: bool


def cited_chunk_ids(answer: str, chunk_ids: list[str]) -> list[str]:
    """Map the [n] source markers the model actually wrote back to chunk ids (1-based,
    in ascending source order, de-duplicated). Out-of-range or null ids are ignored."""
    used = []
    for n in sorted({int(m) for m in re.findall(r"\[(\d+)\]", answer)}):
        if 1 <= n <= len(chunk_ids) and chunk_ids[n - 1]:
            used.append(chunk_ids[n - 1])
    return used


@app.post("/rag/query", response_model=RagQueryResponse, dependencies=[Depends(require_key)])
def rag_query(req: RagQueryRequest, response: Response, correlation_id: str = Depends(correlation_id_dep)):
    response.headers["X-Correlation-Id"] = correlation_id
    request_start = time.monotonic()

    hits = hybrid_retrieve(req.question, req.allowed_document_ids, req.top_k, correlation_id)
    if not hits:
        # Same hallucination guard as before: no context -> don't call the LLM.
        log_stage(correlation_id, "rag_query_total", (time.monotonic() - request_start) * 1000,
                  no_context_found=True)
        return RagQueryResponse(answer="", used_chunk_ids=[], no_context_found=True)

    chunk_ids = [h["chunk_id"] for h in hits]
    # ponytail: content-only context (no title/page — not in vector_store metadata). Java's
    # citations still carry title/page from document_chunks. Upgrade: add them to the ingestion
    # metadata (ChunkIngestionService) or grant Python read on document_chunks (Plan.md §256).
    context = "Context:\n" + "\n\n".join(f"[{i + 1}] {h['content']}" for i, h in enumerate(hits))
    user_message = context + "\n\nQuestion: " + req.question

    t_llm = time.monotonic()
    resp = _groq.chat.completions.create(model=GROQ_MODEL, messages=[
        {"role": "system", "content": SYSTEM_INSTRUCTION},
        {"role": "user", "content": user_message},
    ])
    log_stage(correlation_id, "llm", (time.monotonic() - t_llm) * 1000, model=GROQ_MODEL)

    answer = resp.choices[0].message.content or ""
    log_stage(correlation_id, "rag_query_total", (time.monotonic() - request_start) * 1000,
              no_context_found=False)
    return RagQueryResponse(
        answer=answer,
        used_chunk_ids=cited_chunk_ids(answer, chunk_ids),
        no_context_found=False,
    )


# --- Phase 8: same RAG assembly, but streamed token-by-token. Citations/used_chunk_ids only
# exist once generation finishes, so the stream always ends with one structured "metadata"
# event — Java (and any client) treats that as the signal the answer is complete and citable.

def _sse(payload: dict) -> str:
    return f"data: {json.dumps(payload)}\n\n"


def _stream_rag_query(question: str, allowed_document_ids: list[str], top_k: int, correlation_id: str):
    request_start = time.monotonic()
    hits = hybrid_retrieve(question, allowed_document_ids, top_k, correlation_id)
    if not hits:
        # Same hallucination guard as /rag/query: no context -> don't call the LLM.
        log_stage(correlation_id, "rag_query_stream_total", (time.monotonic() - request_start) * 1000,
                  no_context_found=True)
        yield _sse({"type": "metadata", "used_chunk_ids": [], "no_context_found": True})
        return

    chunk_ids = [h["chunk_id"] for h in hits]
    context = "Context:\n" + "\n\n".join(f"[{i + 1}] {h['content']}" for i, h in enumerate(hits))
    user_message = context + "\n\nQuestion: " + question

    t_llm = time.monotonic()
    first_token_ms = None
    full_answer = []
    stream = _groq.chat.completions.create(model=GROQ_MODEL, stream=True, messages=[
        {"role": "system", "content": SYSTEM_INSTRUCTION},
        {"role": "user", "content": user_message},
    ])
    for chunk in stream:
        delta = chunk.choices[0].delta.content if chunk.choices else None
        if not delta:
            continue
        if first_token_ms is None:
            first_token_ms = (time.monotonic() - t_llm) * 1000
        full_answer.append(delta)
        yield _sse({"type": "token", "content": delta})

    log_stage(correlation_id, "llm", (time.monotonic() - t_llm) * 1000, model=GROQ_MODEL,
              time_to_first_token_ms=round(first_token_ms, 1) if first_token_ms is not None else None)

    answer = "".join(full_answer)
    log_stage(correlation_id, "rag_query_stream_total", (time.monotonic() - request_start) * 1000,
              no_context_found=False)
    yield _sse({
        "type": "metadata",
        "used_chunk_ids": cited_chunk_ids(answer, chunk_ids),
        "no_context_found": False,
    })


@app.post("/rag/query/stream", dependencies=[Depends(require_key)])
def rag_query_stream(req: RagQueryRequest, response: Response, correlation_id: str = Depends(correlation_id_dep)):
    response.headers["X-Correlation-Id"] = correlation_id
    return StreamingResponse(
        _stream_rag_query(req.question, req.allowed_document_ids, req.top_k, correlation_id),
        media_type="text/event-stream",
    )


# --- Phase 5.5: the ingestion WRITE now lives here too. Java persists the authoritative
# DocumentChunk rows (its system of record) and hands us {chunk_id, document_id, content};
# we embed and insert into vector_store. metadata mirrors what /retrieve and /rag/query read.

class IndexChunk(BaseModel):
    chunk_id: str
    document_id: str
    content: str


class IndexRequest(BaseModel):
    chunks: list[IndexChunk]


class IndexResponse(BaseModel):
    indexed: int


# id = chunk_id gives a clean 1:1 chunk↔vector mapping. ponytail: plain INSERT (chunk ids are
# freshly generated per ingestion, so no conflicts today); add ON CONFLICT if a re-index path lands.
_INDEX_SQL = (
    "INSERT INTO vector_store (id, content, metadata, embedding) "
    "VALUES (%(id)s, %(content)s, %(metadata)s, %(embedding)s::vector)"
)


@app.post("/index", response_model=IndexResponse, dependencies=[Depends(require_key)])
def index(req: IndexRequest):
    if not req.chunks:
        return IndexResponse(indexed=0)
    vectors = _embedder.encode([c.content for c in req.chunks])
    _pool.open()  # idempotent
    with _pool.connection() as conn:
        for chunk, vec in zip(req.chunks, vectors):
            conn.execute(_INDEX_SQL, {
                "id": chunk.chunk_id,
                "content": chunk.content,
                "metadata": Json({"chunkId": chunk.chunk_id, "documentId": chunk.document_id}),
                "embedding": "[" + ",".join(map(str, vec.tolist())) + "]",
            })
    return IndexResponse(indexed=len(req.chunks))


# --- Phase 9: OCR + multi-format parsing. Java's PDFBox handles normal PDFs directly (fast, no
# network); this endpoint is only for what PDFBox can't do — scanned/image PDFs (OCR) and
# DOCX/XLSX/PPTX. Response shape mirrors Java's PdfParsingService.ParsedDocument exactly
# ({page_count, pages:[{page_number,text}]}), so DocumentService's downstream chunking pipeline
# doesn't need to know or care which parser produced the pages.

class ParsePage(BaseModel):
    page_number: int
    text: str


class ParseResponse(BaseModel):
    page_count: int
    pages: list[ParsePage]


def _parse_docx(data: bytes) -> ParseResponse:
    # ponytail: DOCX has no native page concept — the whole document is one logical "page".
    # Real pagination would need a layout engine (page size, margins, fonts); not worth it here.
    from docx import Document as DocxDocument
    doc = DocxDocument(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs if p.text)
    return ParseResponse(page_count=1, pages=[ParsePage(page_number=1, text=text)])


def _parse_xlsx(data: bytes) -> ParseResponse:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    pages = []
    for i, sheet in enumerate(wb.worksheets, start=1):
        rows = [" ".join(str(c) for c in row if c is not None)
                for row in sheet.iter_rows(values_only=True)]
        pages.append(ParsePage(page_number=i, text="\n".join(r for r in rows if r)))
    return ParseResponse(page_count=len(pages), pages=pages)


def _parse_pptx(data: bytes) -> ParseResponse:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [shape.text_frame.text for shape in slide.shapes
                 if shape.has_text_frame and shape.text_frame.text]
        pages.append(ParsePage(page_number=i, text="\n".join(texts)))
    return ParseResponse(page_count=len(pages), pages=pages)


def _parse_pdf_ocr(data: bytes) -> ParseResponse:
    """OCR fallback for scanned/image-only PDFs (Java already tried PDFBox and got nothing)."""
    import fitz  # PyMuPDF: rasterize each page, no external binary needed for this half
    import pytesseract
    from PIL import Image
    pages = []
    with fitz.open(stream=data, filetype="pdf") as pdf:
        for i, page in enumerate(pdf, start=1):
            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            pages.append(ParsePage(page_number=i, text=pytesseract.image_to_string(img)))
    return ParseResponse(page_count=len(pages), pages=pages)


_PARSERS = {
    "docx": _parse_docx, "xlsx": _parse_xlsx, "pptx": _parse_pptx, "pdf": _parse_pdf_ocr,
}


@app.post("/parse", response_model=ParseResponse, dependencies=[Depends(require_key)])
async def parse_document(file: UploadFile):
    ext = (file.filename or "").rsplit(".", 1)[-1].lower()
    parser = _PARSERS.get(ext)
    if parser is None:
        raise HTTPException(status_code=422, detail=f"Unsupported format: .{ext}")
    data = await file.read()
    return parser(data)
